#!/usr/bin/env python3
"""better-slides — extract content from a .pptx for conversion.

    python3 scripts/extract-pptx.py <input.pptx> <output_dir>

Writes:
  <output_dir>/content.md              — slide-by-slide text, notes, images (for the agent)
  <output_dir>/extracted-slides.json   — structured data (titles, text, image dims, notes)
  <output_dir>/assets/                 — every embedded image, named slideNN-imgNN.<ext>

Requires python-pptx:  pip install python-pptx

The agent reads content.md, applies the Speaker Playbook to it (find the one
idea this deck is burying), then re-authors it on the better-slides engine —
a performance re-write, not a 1:1 port.
"""
import json
import sys
from pathlib import Path

PICTURE = 13  # MSO_SHAPE_TYPE.PICTURE


def extract(src: Path, outdir: Path) -> list[dict]:
    from pptx import Presentation

    assets = outdir / "assets"
    assets.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(src))
    slides = []

    for n, slide in enumerate(prs.slides, 1):
        data = {"number": n, "title": "", "content": [], "images": [], "notes": ""}
        title_shape = slide.shapes.title

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                if title_shape is not None and shape == title_shape:
                    data["title"] = shape.text.strip()
                else:
                    data["content"].append(shape.text.strip())
            if shape.shape_type == PICTURE:
                image = shape.image
                name = f"slide{n:02d}-img{len(data['images']) + 1:02d}.{image.ext}"
                (assets / name).write_bytes(image.blob)
                data["images"].append(
                    {"path": f"assets/{name}", "width": shape.width, "height": shape.height}
                )

        if slide.has_notes_slide:
            data["notes"] = slide.notes_slide.notes_text_frame.text.strip()
        slides.append(data)

    return slides


def write_markdown(slides: list[dict], outdir: Path, src_name: str) -> None:
    lines = [f"# Extracted from {src_name}", ""]
    for s in slides:
        lines.append(f"## Slide {s['number']}" + (f" — {s['title']}" if s["title"] else ""))
        lines.append("")
        for text in s["content"]:
            for ln in text.splitlines():
                if ln.strip():
                    lines.append(f"- {ln.strip()}")
        for img in s["images"]:
            lines.append(f"- ![image]({img['path']})")
        if s["notes"]:
            lines += ["", f"> speaker notes: {s['notes']}"]
        lines.append("")
    (outdir / "content.md").write_text("\n".join(lines), encoding="utf-8")


def main() -> int:
    if len(sys.argv) != 3:
        print(__doc__)
        return 1
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("✗ python-pptx is required:  pip install python-pptx")
        return 1

    src, outdir = Path(sys.argv[1]), Path(sys.argv[2])
    if not src.is_file():
        print(f"✗ no such file: {src}")
        return 1

    slides = extract(src, outdir)
    write_markdown(slides, outdir, src.name)
    (outdir / "extracted-slides.json").write_text(
        json.dumps(slides, indent=2, ensure_ascii=False), encoding="utf-8"
    )

    total_imgs = sum(len(s["images"]) for s in slides)
    print(f"✓ {len(slides)} slides, {total_imgs} images → {outdir}/content.md")
    for s in slides:
        print(f"  Slide {s['number']}: {s['title'] or '(no title)'} — "
              f"{len(s['content'])} text block(s), {len(s['images'])} image(s)"
              + (", notes" if s["notes"] else ""))
    return 0


if __name__ == "__main__":
    sys.exit(main())
